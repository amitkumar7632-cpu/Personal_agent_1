import os
print("DEBUG: Starting Flask app import...")

import numpy as np
print("DEBUG: Imported numpy")

from flask import Flask, request, render_template
print("DEBUG: Imported Flask")

from pptx import Presentation
print("DEBUG: Imported python-pptx")

from docx import Document
print("DEBUG: Imported python-docx")

import pdfplumber
print("DEBUG: Imported pdfplumber")

import google.generativeai as genai
print("DEBUG: Imported google.generativeai")

from sklearn.neighbors import NearestNeighbors
print("DEBUG: Imported NearestNeighbors")

print("DEBUG: All imports completed successfully!")

app = Flask(__name__)

# -----------------------------
# Lazy Globals
# -----------------------------
all_texts = None
embeddings = None
nn_model = None
genai_configured = False


# -----------------------------
# Extract text from PPT
# -----------------------------
def extract_text_from_ppt(filepath):
    print(f"DEBUG: Extracting PPT: {filepath}")
    texts = []
    prs = Presentation(filepath)
    for i, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        texts.append(f"{os.path.basename(filepath)} - Slide {i}: {para.text.strip()}")
    return texts


# -----------------------------
# Extract text from PDF
# -----------------------------
def extract_text_from_pdf(filepath):
    print(f"DEBUG: Extracting PDF: {filepath}")
    texts = []
    with pdfplumber.open(filepath) as pdf:
        for i, page in enumerate(pdf.pages, start=1):
            text = page.extract_text()
            if text and text.strip():
                texts.append(f"{os.path.basename(filepath)} - Page {i}: {text.strip()}")
    return texts


# -----------------------------
# Ingest folder
# -----------------------------
def ingest_folder(folder_path):
    print(f"DEBUG: Ingesting folder: {folder_path}")
    texts = []
    if not os.path.exists(folder_path):
        print("DEBUG: Folder does NOT exist!")
        return texts

    for filename in os.listdir(folder_path):
        print(f"DEBUG: Reading file: {filename}")

        if filename.startswith("~$"):
            continue

        filepath = os.path.join(folder_path, filename)

        if filename.lower().endswith(".pptx"):
            texts.extend(extract_text_from_ppt(filepath))

        elif filename.lower().endswith(".pdf"):
            texts.extend(extract_text_from_pdf(filepath))

        elif filename.lower().endswith(".docx"):
            doc = Document(filepath)
            for i, para in enumerate(doc.paragraphs, start=1):
                if para.text.strip():
                    texts.append(f"{filename} - Paragraph {i}: {para.text.strip()}")

        elif filename.lower().endswith(".txt"):
            with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
                for i, line in enumerate(f.readlines(), start=1):
                    if line.strip():
                        texts.append(f"{filename} - Line {i}: {line.strip()}")

    print("DEBUG: Finished ingesting folder")
    return texts


# -----------------------------
# Highlight helper
# -----------------------------
def highlight_terms(text, query):
    for term in query.split():
        text = text.replace(term, f"<mark>{term}</mark>")
        text = text.replace(term.capitalize(), f"<mark>{term.capitalize()}</mark>")
    return text


# -----------------------------
# Home Route
# -----------------------------
@app.route("/", methods=["GET", "POST"])
def home():
    print("DEBUG: Home route triggered")
    global all_texts, embeddings, nn_model, genai_configured

    answer = ""
    context = ""

    if request.method == "POST":
        print("DEBUG: POST request received")

        query = request.form["query"]
        print(f"DEBUG: Query received: {query}")

        # Load documents
        if all_texts is None:
            print("DEBUG: Loading ALL_Docs...")
            all_texts = ingest_folder("ALL_Docs")

        if not all_texts:
            print("DEBUG: No documents found!")
            return "No documents found in ALL_Docs folder."

        # Configure Gemini once
        if not genai_configured:
            print("DEBUG: Configuring Gemini API...")
            genai.configure(api_key=os.getenv("GEMINI_API_KEY"))
            genai_configured = True

        # Build embeddings using Gemini
        if embeddings is None:
            print("DEBUG: Generating embeddings for documents...")
            embed_response = genai.embed_content(
                model="models/text-embedding-004",
                content=all_texts
            )
            embeddings = np.array(embed_response["embedding"])
            print("DEBUG: Embeddings generated!")

            nn_model = NearestNeighbors(n_neighbors=5, metric="cosine")
            nn_model.fit(embeddings)
            print("DEBUG: NearestNeighbors model built!")

        # Query embedding
        print("DEBUG: Generating query embedding...")
        query_embed = genai.embed_content(
            model="models/text-embedding-004",
            content=query
        )["embedding"]

        print("DEBUG: Running nearest neighbor search...")
        distances, indices = nn_model.kneighbors([query_embed])

        context = "\n".join([highlight_terms(all_texts[idx], query) for idx in indices[0]])

        print("DEBUG: Sending prompt to Gemini...")
        response = genai.generate_text(
            model="models/gemini-2.5-flash",
            prompt=(
                f"Answer the question strictly using the context below.\n\n"
                f"Context:\n{context}\n\n"
                f"Question: {query}\n\n"
                f"Give a direct answer based only on the context."
            )
        )

        print("DEBUG: Gemini response received")
        answer = response.result

    return render_template("index.html", answer=answer, context=context)


# -----------------------------
# Run locally
# -----------------------------
if __name__ == "__main__":
    print("DEBUG: Starting Flask server...")
    port = int(os.getenv("PORT", 5000))
    app.run(host="0.0.0.0", port=port)

from flask import Flask, request, render_template
from sentence_transformers import SentenceTransformer
import faiss
import numpy as np
from pptx import Presentation
from google import genai

# Load PPT file
ppt_path = r"D:\HAVELLS\17102025_ICE MAKER & GRILL ASSEMBLY_SV REVISED.pptx"
prs = Presentation(ppt_path)
ppt_text = []
for i, slide in enumerate(prs.slides, start=1):
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            ppt_text.append(f"Slide {i}: {shape.text}")

# Embeddings + FAISS
model = SentenceTransformer("all-MiniLM-L6-v2")
embeddings = model.encode(ppt_text)
dimension = embeddings.shape[1]
index = faiss.IndexFlatL2(dimension)
index.add(np.array(embeddings))

# Gemini client
client = genai.Client(api_key="AIzaSyAqXiJcrigh5bxeiWPnspnLcOWyCRcGqVE")

app = Flask(__name__)

@app.route("/", methods=["GET", "POST"])
def home():
    answer = ""
    context = ""
    if request.method == "POST":
        query = request.form["query"]
        query_embedding = model.encode([query])
        D, I = index.search(np.array(query_embedding), k=5)
        context = "\n".join([ppt_text[idx] for idx in I[0]])

        response = client.models.generate_content(
            model="models/gemini-2.5-flash",
            contents=f"Use this context:\n{context}\n\nQuestion: {query}"
        )
        answer = response.text

    return render_template("index.html", answer=answer, context=context)

if __name__ == "__main__":
    app.run(debug=True)
from pptx import Presentation

ppt_path = r"D:\HAVELLS\17102025_ICE MAKER & GRILL ASSEMBLY_SV REVISED.pptx"
prs = Presentation(ppt_path)

# Save extracted text into a file
with open("ppt_text_output.txt", "w", encoding="utf-8") as f:
    for i, slide in enumerate(prs.slides, start=1):
        f.write(f"\n--- Slide {i} ---\n")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                f.write(shape.text + "\n")

print("Text extracted and saved to ppt_text_output.txt")

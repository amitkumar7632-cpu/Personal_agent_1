from pptx import Presentation
from sentence_transformers import SentenceTransformer
import faiss
import numpy as np
from google import genai

# -----------------------------
# Step 1: Extract text from PPTX
# -----------------------------
ppt_path = r"D:\HAVELLS\17102025_ICE MAKER & GRILL ASSEMBLY_SV REVISED.pptx"
prs = Presentation(ppt_path)

ppt_text = []
for i, slide in enumerate(prs.slides, start=1):
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            ppt_text.append(f"Slide {i}: {shape.text}")

print(f"Extracted {len(ppt_text)} text blocks from PPTX")

# -----------------------------
# Step 2: Embed text with SentenceTransformer
# -----------------------------
model = SentenceTransformer("all-MiniLM-L6-v2")
embeddings = model.encode(ppt_text)

# -----------------------------
# Step 3: Store embeddings in FAISS
# -----------------------------
dimension = embeddings.shape[1]
index = faiss.IndexFlatL2(dimension)
index.add(np.array(embeddings))

print("Embeddings stored in FAISS index!")

# -----------------------------
# Step 4: Initialize Gemini client
# -----------------------------
client = genai.Client(api_key="AIzaSyAqXiJcrigh5bxeiWPnspnLcOWyCRcGqVE")

# -----------------------------
# Step 5: Interactive Q&A loop
# -----------------------------
while True:
    query = input("\nAsk your question (or type 'exit'): ")
    if query.lower() == "exit":
        print("Exiting agent. Goodbye!")
        break

    # Embed query
    query_embedding = model.encode([query])
    D, I = index.search(np.array(query_embedding), k=5)  # retrieve top 5 matches

    # Collect context
    context = "\n".join([ppt_text[idx] for idx in I[0]])

    # Ask Gemini with context
    response = client.models.generate_content(
        model="models/gemini-2.5-flash",   # or "models/gemini-pro-latest"
        contents=f"Use this context:\n{context}\n\nQuestion: {query}"
    )

    print("\nGemini Answer:\n", response.text)